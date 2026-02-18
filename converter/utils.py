from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
import os
import tempfile
import json
import zipfile
import base64
import subprocess
import shutil
import uuid
from PIL import Image


def pdf_to_pptx(pdf_path, output_path, progress_callback=None):

    # Convert PDF pages to images
    with tempfile.TemporaryDirectory() as temp_dir:
        images = convert_from_path(pdf_path, dpi=150)

        # Create PowerPoint presentation
        prs = Presentation()

        # Set slide dimensions (16:9 widescreen aspect ratio)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

        total_pages = len(images)

        for i, image in enumerate(images):
            # Add blank slide
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)

            # Save image temporarily
            img_path = os.path.join(temp_dir, f'page_{i}.png')
            image.save(img_path, 'PNG')

            # Add image to slide (fit to slide)
            left = Inches(0)
            top = Inches(0)
            slide.shapes.add_picture(img_path, left, top,
                                     width=prs.slide_width,
                                     height=prs.slide_height)

            # Report progress
            if progress_callback:
                progress_callback(i + 1, total_pages)

        # Save presentation
        prs.save(output_path)

    return output_path


def pptx_to_images(pptx_path, output_dir, progress_callback=None):
    """
    Convert PPTX slides to images using LibreOffice and pdf2image

    Args:
        pptx_path: Path to PPTX file
        output_dir: Directory where images will be saved
        progress_callback: Optional callback for progress updates

    Returns:
        List of image file paths
    """
    image_paths = []

    with tempfile.TemporaryDirectory() as temp_pdf_dir:
        # Step 1: Convert PPTX to PDF using LibreOffice
        try:
            result = subprocess.run(
                ['soffice', '--headless', '--convert-to', 'pdf', '--outdir', temp_pdf_dir, pptx_path],
                capture_output=True,
                text=True,
                timeout=300  # 5 minute timeout
            )

            if result.returncode != 0:
                raise Exception(f"LibreOffice conversion failed: {result.stderr}")

            # Find the generated PDF file
            pdf_filename = os.path.splitext(os.path.basename(pptx_path))[0] + '.pdf'
            pdf_path = os.path.join(temp_pdf_dir, pdf_filename)

            if not os.path.exists(pdf_path):
                raise Exception(f"PDF file not generated: {pdf_path}")

        except subprocess.TimeoutExpired:
            raise Exception("PPTX conversion timeout (exceeded 5 minutes)")
        except FileNotFoundError:
            raise Exception("LibreOffice (soffice) not found. Please ensure LibreOffice is installed.")
        except Exception as e:
            raise Exception(f"Error converting PPTX to PDF: {str(e)}")

        # Step 2: Convert PDF to images using pdf2image
        try:
            images = convert_from_path(pdf_path, dpi=150)
            total_slides = len(images)

            for i, image in enumerate(images):
                # Save each slide as an image
                img_filename = f'pptx_slide_{i}.png'
                img_path = os.path.join(output_dir, img_filename)
                image.save(img_path, 'PNG')
                image_paths.append(img_path)

                # Report progress
                if progress_callback:
                    progress_callback(i + 1, total_slides, 'Converting PPTX slides')

        except Exception as e:
            raise Exception(f"Error converting PDF to images: {str(e)}")

    return image_paths


def resize_to_aspect_ratio(img, target_ratio=2.0):
    """
    Resize image to target aspect ratio (width/height)

    Args:
        img: PIL Image object
        target_ratio: Target aspect ratio (default 2.0 for 2:1)

    Returns:
        PIL Image object with target aspect ratio
    """
    original_width, original_height = img.size
    current_ratio = original_width / original_height

    if abs(current_ratio - target_ratio) < 0.01:
        # Already close to target ratio
        return img

    # Calculate new dimensions to fit target ratio
    if current_ratio > target_ratio:
        # Image is too wide, crop width
        new_width = int(original_height * target_ratio)
        new_height = original_height
        # Center crop
        left = (original_width - new_width) // 2
        top = 0
        right = left + new_width
        bottom = original_height
    else:
        # Image is too tall, crop height
        new_width = original_width
        new_height = int(original_width / target_ratio)
        # Center crop
        left = 0
        top = (original_height - new_height) // 2
        right = original_width
        bottom = top + new_height

    # Crop to target aspect ratio
    cropped_img = img.crop((left, top, right, bottom))

    # Optionally resize to a standard size (e.g., 1600x800 for 2:1)
    # This ensures consistent sizing across all images
    standard_width = 1600
    standard_height = int(standard_width / target_ratio)
    resized_img = cropped_img.resize((standard_width, standard_height), Image.LANCZOS)

    return resized_img


def images_to_h5p(image_paths, output_path, content_type='presentation', alignment='middle', progress_callback=None):
    """
    Convert images to H5P format (Course Presentation or Interactive Book)

    Args:
        image_paths: List of paths to image files
        output_path: Path where the .h5p file will be saved
        content_type: 'presentation' or 'interactive-book'
        alignment: 'left', 'middle', 'right', or 'fullscreen'
        progress_callback: Optional callback function for progress updates
    """

    with tempfile.TemporaryDirectory() as temp_dir:
        # Create H5P directory structure
        h5p_dir = os.path.join(temp_dir, 'h5p_content')
        os.makedirs(h5p_dir, exist_ok=True)

        content_dir = os.path.join(h5p_dir, 'content')
        os.makedirs(content_dir, exist_ok=True)

        images_dir = os.path.join(content_dir, 'images')
        os.makedirs(images_dir, exist_ok=True)

        total_images = len(image_paths)

        # Process and copy images
        image_files = []
        for i, img_path in enumerate(image_paths):
            # Open image and convert to 2:1 aspect ratio
            img = Image.open(img_path)

            # Convert to RGB if necessary (for PNG with transparency, etc.)
            if img.mode in ('RGBA', 'LA', 'P'):
                # Create white background
                background = Image.new('RGB', img.size, (255, 255, 255))
                if img.mode == 'P':
                    img = img.convert('RGBA')
                background.paste(img, mask=img.split()[-1] if img.mode in ('RGBA', 'LA') else None)
                img = background
            elif img.mode != 'RGB':
                img = img.convert('RGB')

            # Resize to 2:1 aspect ratio (width = 2 × height)
            img = resize_to_aspect_ratio(img, target_ratio=2.0)

            img_filename = f'image_{i}.png'  # Always save as PNG for consistency
            img_save_path = os.path.join(images_dir, img_filename)

            # Save image
            img.save(img_save_path, 'PNG')
            image_files.append({
                'filename': img_filename,
                'path': f'images/{img_filename}',
                'width': img.width,
                'height': img.height
            })

            if progress_callback:
                progress_callback(i + 1, total_images, 'Processing images')

        # Create content.json based on content type
        if content_type == 'presentation':
            content_json = create_presentation_content(image_files, alignment)
            library = 'H5P.CoursePresentation'
            major_version = 1
            minor_version = 25
            dependencies = [
                {"machineName": "H5P.CoursePresentation", "majorVersion": 1, "minorVersion": 25},
                {"machineName": "H5P.Image", "majorVersion": 1, "minorVersion": 1}
            ]
        else:
            content_json = create_interactive_book_content(image_files, alignment)
            library = 'H5P.InteractiveBook'
            major_version = 1
            minor_version = 8
            dependencies = [
                {"machineName": "H5P.InteractiveBook", "majorVersion": 1, "minorVersion": 8},
                {"machineName": "H5P.Column", "majorVersion": 1, "minorVersion": 16},
                {"machineName": "H5P.AdvancedText", "majorVersion": 1, "minorVersion": 1}
            ]

        # Write content.json
        with open(os.path.join(content_dir, 'content.json'), 'w', encoding='utf-8') as f:
            json.dump(content_json, f, indent=2)

        # Create h5p.json
        h5p_json = {
            "title": "Image Collection",
            "language": "en",
            "mainLibrary": library,
            "embedTypes": ["div"],
            "license": "U",
            "preloadedDependencies": dependencies
        }

        with open(os.path.join(h5p_dir, 'h5p.json'), 'w', encoding='utf-8') as f:
            json.dump(h5p_json, f, indent=2)

        # Create the H5P package (zip file)
        with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            # Add h5p.json
            zipf.write(os.path.join(h5p_dir, 'h5p.json'), 'h5p.json')

            # Add content directory
            for root, dirs, files in os.walk(content_dir):
                for file in files:
                    file_path = os.path.join(root, file)
                    arcname = os.path.relpath(file_path, h5p_dir)
                    zipf.write(file_path, arcname)

        if progress_callback:
            progress_callback(total_images, total_images, 'H5P package created')

    return output_path


def create_presentation_content(image_files, alignment):
    """Create H5P Course Presentation content structure"""

    # x = left, y = top
    alignment_map = {
        'left': {'x': 0, 'width': 70, 'height': 70, 'y': 5},
        'middle': {'x': 5, 'width': 70, 'height': 70, 'y': 5},
        'right': {'x': 30, 'width': 70, 'height': 70, 'y': 5},
        'fullscreen': {'x': 0, 'width': 100, 'height': 100, 'y': 0}
    }

    pos = alignment_map.get(alignment, alignment_map['middle'])

    slides = []
    for img_data in image_files:
        slide = {
            "elements": [
                {
                    "x": pos['x'],
                    "y": pos['y'],
                    "width": pos['width'],
                    "height": pos['height'],
                    "action": {
                        "library": "H5P.Image 1.1",
                        "params": {
                            "contentName": "Image",
                            "file": {
                                "path": img_data['path'],
                                "mime": f"image/{img_data['filename'].split('.')[-1]}",
                                "width": img_data['width'],
                                "height": img_data['height']
                            },
                            "alt": "Image"
                        },
                        "subContentId": str(uuid.uuid4())
                    }
                }
            ],
            "keywords": []
        }
        slides.append(slide)

    return {
        "presentation": {
            "slides": slides
        },
        "override": {
            "activeSurface": False,
            "hideSummarySlide": True,
            "summarySlideImage": None
        }
    }


def create_interactive_book_content(image_files, alignment):
    """Create H5P Interactive Book content structure"""

    # Map alignment to CSS text-align values
    alignment_css = {
        'left': 'left',
        'middle': 'center',
        'right': 'right',
        'fullscreen': 'center'
    }

    # Map alignment to image width
    alignment_width = {
        'left': '50%',
        'middle': '70%',
        'right': '50%',
        'fullscreen': '100%'
    }

    css_align = alignment_css.get(alignment, 'center')
    img_width = alignment_width.get(alignment, '70%')

    chapters = []
    for i, img_data in enumerate(image_files):
        # Create chapter with Column layout containing AdvancedText with image
        chapter = {
            "title": f"Image {i + 1}",
            "content": {
                "library": "H5P.Column 1.16",
                "params": {
                    "content": [
                        {
                            "content": {
                                "library": "H5P.AdvancedText 1.1",
                                "params": {
                                    "text": f'<p style="text-align: {css_align};"><img src="images/{img_data["filename"]}" alt="Image {i + 1}" style="width: {img_width}; height: auto; max-width: 100%;" /></p>'
                                },
                                "subContentId": str(uuid.uuid4()),
                                "metadata": {
                                    "contentType": "Text",
                                    "license": "U",
                                    "title": f"Image {i + 1}"
                                }
                            },
                            "useSeparator": "auto"
                        }
                    ]
                },
                "subContentId": str(uuid.uuid4()),
                "metadata": {
                    "contentType": "Column",
                    "license": "U",
                    "title": f"Chapter {i + 1}"
                }
            }
        }
        chapters.append(chapter)

    return {
        "bookCover": {
            "coverDescription": "Image collection as interactive book",
            "coverMedium": None,
            "coverAltText": ""
        },
        "chapters": chapters,
        "behaviour": {
            "defaultTableOfContents": True,
            "progressIndicators": True,
            "progressAuto": False,
            "displayTOC": True
        },
        "read": "Read",
        "displayTOC": "Display table of contents",
        "hideTOC": "Hide table of contents",
        "nextPage": "Next page",
        "previousPage": "Previous page",
        "navigationLabel": "Book navigation",
        "titleLabel": "Book title",
        "pagesLabel": "Pages",
        "pageLabel": "Page",
        "readspeakerProgress": "@pages of @total pages"
    }


def pdf_to_images_zip(pdf_path, output_zip_path, progress_callback=None):
    """
    Convert PDF pages to PNG images and package them in a ZIP file.

    Args:
        pdf_path: Path to the PDF file
        output_zip_path: Path where the .zip file will be saved
        progress_callback: Optional callback(current, total)

    Returns:
        output_zip_path
    """
    with tempfile.TemporaryDirectory() as temp_dir:
        images = convert_from_path(pdf_path, dpi=150)
        total_pages = len(images)

        png_paths = []
        for i, image in enumerate(images):
            img_path = os.path.join(temp_dir, f'page_{i + 1}.png')
            image.save(img_path, 'PNG')
            png_paths.append(img_path)

            if progress_callback:
                progress_callback(i + 1, total_pages)

        with zipfile.ZipFile(output_zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for path in png_paths:
                zipf.write(path, os.path.basename(path))

    return output_zip_path


IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png', '.webp', '.gif', '.bmp', '.tiff', '.tif'}
VIDEO_EXTENSIONS = {'.mp4', '.avi', '.mov', '.mkv', '.webm', '.wmv', '.flv'}
AUDIO_EXTENSIONS = {'.mp3', '.wav', '.flac', '.aac', '.ogg', '.m4a', '.wma'}
PDF_EXTENSIONS = {'.pdf'}


def get_file_type(filename):
    ext = os.path.splitext(filename)[1].lower()
    if ext in IMAGE_EXTENSIONS:
        return 'image'
    if ext in VIDEO_EXTENSIONS:
        return 'video'
    if ext in AUDIO_EXTENSIONS:
        return 'audio'
    if ext in PDF_EXTENSIONS:
        return 'pdf'
    return None


def compress_image(input_path, output_path):
    img = Image.open(input_path)

    # Resize if any dimension exceeds 4096px
    max_dim = 4096
    if img.width > max_dim or img.height > max_dim:
        ratio = min(max_dim / img.width, max_dim / img.height)
        new_size = (int(img.width * ratio), int(img.height * ratio))
        img = img.resize(new_size, Image.LANCZOS)

    # Handle transparency: keep as PNG, otherwise save as JPEG
    if img.mode in ('RGBA', 'LA', 'P'):
        img.save(output_path, 'PNG', optimize=True)
    else:
        if img.mode != 'RGB':
            img = img.convert('RGB')
        img.save(output_path, 'JPEG', quality=70, optimize=True)

    return output_path


def compress_video(input_path, output_path):
    result = subprocess.run(
        [
            'ffmpeg', '-i', input_path,
            '-vcodec', 'libx264', '-crf', '28',
            '-acodec', 'aac', '-b:a', '128k',
            '-vf', 'scale=min(iw\\,1920):min(ih\\,1080):force_original_aspect_ratio=decrease',
            '-movflags', '+faststart',
            '-y', output_path
        ],
        capture_output=True, text=True, timeout=600
    )
    if result.returncode != 0:
        raise Exception(f"FFmpeg video compression failed: {result.stderr}")
    return output_path


def compress_audio(input_path, output_path):
    result = subprocess.run(
        [
            'ffmpeg', '-i', input_path,
            '-acodec', 'libmp3lame', '-b:a', '128k',
            '-ar', '44100',
            '-y', output_path
        ],
        capture_output=True, text=True, timeout=300
    )
    if result.returncode != 0:
        raise Exception(f"FFmpeg audio compression failed: {result.stderr}")
    return output_path


def compress_pdf(input_path, output_path):
    result = subprocess.run(
        [
            'gs', '-sDEVICE=pdfwrite',
            '-dCompatibilityLevel=1.4',
            '-dPDFSETTINGS=/ebook',
            '-dNOPAUSE', '-dQUIET', '-dBATCH',
            f'-sOutputFile={output_path}',
            input_path
        ],
        capture_output=True, text=True, timeout=600
    )
    if result.returncode != 0:
        raise Exception(f"Ghostscript PDF-Komprimierung fehlgeschlagen: {result.stderr}")
    return output_path


def compress_files(file_paths, output_zip_path, progress_callback=None):
    total = len(file_paths)

    with tempfile.TemporaryDirectory() as temp_dir:
        compressed_paths = []

        for i, file_path in enumerate(file_paths):
            filename = os.path.basename(file_path)
            file_type = get_file_type(filename)
            name, ext = os.path.splitext(filename)

            if file_type == 'image':
                # Output as .jpg unless it has transparency (handled inside compress_image)
                has_alpha = ext.lower() in ('.png', '.gif', '.tiff', '.tif', '.webp')
                if has_alpha:
                    try:
                        img = Image.open(file_path)
                        has_alpha = img.mode in ('RGBA', 'LA', 'P')
                        img.close()
                    except:
                        has_alpha = False
                out_ext = '.png' if has_alpha else '.jpg'
                out_path = os.path.join(temp_dir, name + out_ext)
                compress_image(file_path, out_path)
                compressed_paths.append(out_path)

            elif file_type == 'video':
                out_path = os.path.join(temp_dir, name + '.mp4')
                compress_video(file_path, out_path)
                compressed_paths.append(out_path)

            elif file_type == 'audio':
                out_path = os.path.join(temp_dir, name + '.mp3')
                compress_audio(file_path, out_path)
                compressed_paths.append(out_path)

            elif file_type == 'pdf':
                out_path = os.path.join(temp_dir, name + '.pdf')
                compress_pdf(file_path, out_path)
                compressed_paths.append(out_path)

            if progress_callback:
                progress_callback(i + 1, total)

        # Create zip of all compressed files
        with zipfile.ZipFile(output_zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for path in compressed_paths:
                zipf.write(path, os.path.basename(path))

    return output_zip_path